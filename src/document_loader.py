import os
from pathlib import Path
from typing import List
from llama_index.core import Document

class DocumentLoader:
    """Loads documents from various file types in the data/raw directory."""
    
    def __init__(self, data_dir: str = "data/raw"):
        self.data_dir = Path(data_dir)
        
    def load_text_file(self, file_path: Path) -> Document:
        """Load plain text files (.txt, .py, .sql, .md, etc.)"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        metadata = {
            "filename": file_path.name,
            "file_type": file_path.suffix,
            "file_path": str(file_path)
        }
        return Document(text=content, metadata=metadata)
    
    def load_pdf(self, file_path: Path) -> Document:
        """Load PDF files"""
        from pypdf import PdfReader
        
        reader = PdfReader(file_path)
        content = ""
        for page in reader.pages:
            content += page.extract_text() + "\n"
        
        metadata = {
            "filename": file_path.name,
            "file_type": ".pdf",
            "file_path": str(file_path),
            "num_pages": len(reader.pages)
        }
        return Document(text=content, metadata=metadata)
    
    def load_docx(self, file_path: Path) -> Document:
        """Load Word documents"""
        from docx import Document as DocxDocument
        
        doc = DocxDocument(file_path)
        content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        
        metadata = {
            "filename": file_path.name,
            "file_type": ".docx",
            "file_path": str(file_path)
        }
        return Document(text=content, metadata=metadata)
    
    def load_xlsx(self, file_path: Path) -> Document:
        """Load Excel files"""
        import pandas as pd
        
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        content = f"Excel file: {file_path.name}\n\n"
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            content += f"Sheet: {sheet_name}\n"
            content += df.to_string() + "\n\n"
        
        metadata = {
            "filename": file_path.name,
            "file_type": ".xlsx",
            "file_path": str(file_path),
            "num_sheets": len(excel_file.sheet_names)
        }
        return Document(text=content, metadata=metadata)
    
    def load_pptx(self, file_path: Path) -> Document:
        """Load PowerPoint files"""
        from pptx import Presentation
        
        prs = Presentation(file_path)
        content = f"PowerPoint: {file_path.name}\n\n"
        
        for i, slide in enumerate(prs.slides, 1):
            content += f"Slide {i}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
            content += "\n"
        
        metadata = {
            "filename": file_path.name,
            "file_type": ".pptx",
            "file_path": str(file_path),
            "num_slides": len(prs.slides)
        }
        return Document(text=content, metadata=metadata)
    
    def load_all_documents(self) -> List[Document]:
        """Load all supported documents from data directory"""
        documents = []
        
        # Define file type handlers
        handlers = {
            '.txt': self.load_text_file,
            '.py': self.load_text_file,
            '.sql': self.load_text_file,
            '.md': self.load_text_file,
            '.json': self.load_text_file,
            '.csv': self.load_text_file,
            '.pdf': self.load_pdf,
            '.docx': self.load_docx,
            '.doc': self.load_docx,
            '.xlsx': self.load_xlsx,
            '.xls': self.load_xlsx,
            '.pptx': self.load_pptx
        }
        
        # Walk through all files in data directory
        for file_path in self.data_dir.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in handlers:
                try:
                    print(f"Loading: {file_path.name}")
                    handler = handlers[file_path.suffix.lower()]
                    doc = handler(file_path)
                    documents.append(doc)
                except Exception as e:
                    print(f"Error loading {file_path.name}: {str(e)}")
        
        print(f"\nSuccessfully loaded {len(documents)} documents")
        return documents